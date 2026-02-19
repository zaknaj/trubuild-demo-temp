import os
import fitz
import py7zr

fitz.TOOLS.mupdf_display_errors(False)
fitz.TOOLS.mupdf_display_warnings(False)
import io
import time
import docx
import base64
import shutil
import hashlib
import zipfile
import tempfile
import pypandoc
import warnings
import threading
import pandas as pd
from pathlib import Path
from contextlib import contextmanager
from PIL import Image as PILImage, ImageFile
from typing import List, Dict, Any, Optional, Tuple
from concurrent.futures import ThreadPoolExecutor, as_completed
from utils.core.log import get_logger, set_logger, pid_tool_logger

# from utils.bucket import download_single_file  # TODO

from google.genai import types
from google.genai.types import Part

"""
pip install python-docx openpyxl pypandoc pandas Pillow pymupdf
sudo for pypandoc
"""

MAX_INLINE_PAYLOAD_BYTES = 20 * 1024 * 1024
MAX_IMAGES_PER_PAGE = 3
_COUNT_MIN_PIXELS = 45_000 #only used to decide the "too many images" threshold, NOT to drop images.
warnings.simplefilter("ignore", PILImage.DecompressionBombWarning)
PILImage.MAX_IMAGE_PIXELS = None
ImageFile.LOAD_TRUNCATED_IMAGES = True

SUPPORTED_DOCS = {
    ".pdf",
    ".docx",
    ".doc",
    ".txt",
    ".xlsx",
    ".png",
    ".jpg",
    ".jpeg",
    ".gif",
    ".bmp",
    ".pptx",
}
ARCHIVES = {".zip", ".7z"}
_IMAGE_PROCESS_WORKERS = 4

# Global lock: PyMuPDF is NOT thread-safe. Never call fitz.open / page.get_* from
# multiple threads in the same process without holding this.
_pymupdf_lock = threading.Lock()


def _normalize_for_jpeg(im: PILImage.Image) -> PILImage.Image:
    """
    Normalize images before JPEG encode.
    Handles palette images with transparency safely to avoid PIL warnings and
    preserve visual appearance.
    """
    if im.mode == "P":
        transparency = im.info.get("transparency")
        if transparency is not None:
            rgba = im.convert("RGBA")
            bg = PILImage.new("RGBA", rgba.size, (255, 255, 255, 255))
            bg.alpha_composite(rgba)
            return bg.convert("RGB")
        return im.convert("RGB")

    if im.mode == "RGBA":
        bg = PILImage.new("RGBA", im.size, (255, 255, 255, 255))
        bg.alpha_composite(im)
        return bg.convert("RGB")

    if im.mode not in ("RGB", "L"):
        return im.convert("RGB")

    return im if im.mode == "RGB" else im.convert("RGB")


def _px_from_info(info: dict) -> int:
    """Pixel area without decoding bytes."""
    w = info.get("width") or 0
    h = info.get("height") or 0
    try:
        return int(w) * int(h)
    except Exception:
        return 0

def prepare_image_part_bytes(
    image_bytes: bytes,
    *,
    max_part_size_bytes: int = 2 * 1024 * 1024,
    strict_target_kb: int = 200,
    max_dimension: int = 1600,
    _logger=None,
) -> Tuple[Optional[bytes], Optional[str]]:
    """
    Normalize standalone image -> JPEG and enforce max_part_size_bytes.
    Returns (bytes, mime_type) or (None, None) if cannot be shrunk.
    """
    logger = _logger or get_logger()
    try:
        # Validate + normalize to RGB
        with PILImage.open(io.BytesIO(image_bytes)) as im:
            im = _normalize_for_jpeg(im)

            buf = io.BytesIO()
            im.save(buf, format="JPEG", quality=85, optimize=True)
            jpeg_bytes = buf.getvalue()

        # Enforce same cap as other media Parts
        capped = resize_image(
            jpeg_bytes,
            max_part_size_bytes=max_part_size_bytes,
            strict_target_kb=strict_target_kb,
            max_dimension=max_dimension,
            _logger=logger,
        )
        if not capped:
            return None, None

        return capped, "image/jpeg"

    except Exception as e:
        logger.debug(f"prepare_image_part_bytes failed: {e}")
        return None, None

def _collapse_layered_images_to_screenshot_meta(
    *,
    page,
    filename: str,
    page_no: int,
    image_infos: List[dict],
    content: List[Dict[str, Any]],
    errors: List[Dict[str, Any]],
    area_similarity_threshold: float = 0.8,
) -> bool:

    logger = get_logger()

    pxs = [(_px_from_info(i), i) for i in (image_infos or [])]
    pxs = [(px, i) for (px, i) in pxs if px > 0]

    if len(pxs) < 2:
        return False

    pxs.sort(key=lambda t: t[0], reverse=True)
    max_px = pxs[0][0]
    if max_px <= 0:
        return False

    layered = [t for t in pxs if t[0] >= area_similarity_threshold * max_px]

    if len(layered) >= 2:
        try:
            page_png = render_pdf_page_as_image(page)
            content.append(
                {
                    "type": "full_page_image",
                    "page": page_no,
                    "content": base64.b64encode(page_png).decode("utf-8"),
                    "format": "png",
                    "filename": filename,
                }
            )
            logger.debug(
                "Page %d: detected %d layered large images via metadata; replaced with full_page_image.",
                page_no,
                len(layered),
            )
            return True
        except Exception as e:
            tech = f"Page {page_no}: full-page render failed in layered-collapse: {e}"
            logger.debug(tech)
            errors.append(_wrap_error(filename, tech, user_friendly=True))
            return False

    return False


def estimate_total_parts_size(parts: List[Part]) -> int:
    total_size = 0
    for p in parts:
        if getattr(p, "text", None) is not None:
            total_size += len(p.text.encode("utf-8"))
        elif getattr(p, "inline_data", None):
            total_size += len(p.inline_data.data)
    return total_size

def _image_pixel_count(b: bytes) -> int:
    try:
        with PILImage.open(io.BytesIO(b)) as im:
            return (im.width or 0) * (im.height or 0)
    except Exception:
        return 0

def _collapse_layered_images_to_screenshot(
    *,
    page,
    filename: str,
    page_no: int,
    valid_images: List[Tuple[bytes, str]],
    content: List[Dict[str, Any]],
    errors: List[Dict[str, Any]],
    area_similarity_threshold: float = 0.8,
) -> bool:
    """
    Heuristic: if a page has multiple large images whose pixel areas are very
    similar (typical for scanned PDFs with layered foreground/background), we
    treat them as layers of the same page and keep ONLY a full-page screenshot.

    Returns True if we handled the page (added a full_page_image and caller
    should skip normal per-image handling), False otherwise.
    """
    logger = get_logger()

    # Compute pixel counts
    meta = []
    for b, ext in valid_images:
        px = _image_pixel_count(b)
        if px > 0:
            meta.append((b, ext, px))

    if len(meta) < 2:
        return False  # nothing to collapse

    # Sort by area, largest first
    meta.sort(key=lambda t: t[2], reverse=True)
    max_px = meta[0][2]
    if max_px <= 0:
        return False

    # Images whose area is "close" to the largest are considered layers
    layered = [m for m in meta if m[2] >= area_similarity_threshold * max_px]

    # If we have 2+ such "almost full-page" images, assume layers and screenshot
    if len(layered) >= 2:
        try:
            page_png = render_pdf_page_as_image(page)
            content.append(
                {
                    "type": "full_page_image",
                    "page": page_no,
                    "content": base64.b64encode(page_png).decode("utf-8"),
                    "format": "png",
                    "filename": filename,
                }
            )
            logger.debug(
                "Page %d: detected %d layered large images; replaced with full_page_image.",
                page_no,
                len(layered),
            )
            return True
        except Exception as e:
            tech = f"Page {page_no}: full-page render failed in layered-collapse: {e}"
            logger.debug(tech)
            errors.append(_wrap_error(filename, tech, user_friendly=True))
            return False

    return False

_stderr_lock = threading.Lock()
@contextmanager
def muted_c_stderr():
    """
    Temporarily redirect C-level stderr (fd 2) to /dev/null.

    - Thread-safe via a global lock.
    - Guaranteed to restore stderr exactly once.
    - Use as:  with muted_c_stderr():  <call noisy C code>
    """
    devnull = None
    old_stderr_fd = None

    # Serialize access so only one thread tampers with fd 2 at a time
    with _stderr_lock:
        try:
            devnull = os.open(os.devnull, os.O_WRONLY)
            old_stderr_fd = os.dup(2)
            os.dup2(devnull, 2)
            yield
        finally:
            # Best-effort restore; don't crash if something went weird
            try:
                if old_stderr_fd is not None:
                    os.dup2(old_stderr_fd, 2)
            finally:
                if old_stderr_fd is not None:
                    try:
                        os.close(old_stderr_fd)
                    except OSError:
                        pass
                if devnull is not None:
                    try:
                        os.close(devnull)
                    except OSError:
                        pass


def _restore_c_stderr(old_stderr_fd):
    """Restore C-level stderr from saved fd."""
    os.dup2(old_stderr_fd, 2)
    os.close(old_stderr_fd)


def _wrap_error(
    file_name: str, technical_msg: str, *, user_friendly: bool = True
) -> Dict[str, Any]:
    """
    Return a dict with BOTH messages.

    Keys
    ----
    file         - filename the error relates to
    technical    - full stack / exception text (always present)
    errors       - list with the message you want to show to the user
    """
    friendly_msg = (
        "We couldn't process this document/page (s). "
        "The file/page may be damaged or in an unsupported format. "
        "Please check the file and try uploading it again."
    )

    return {
        "file": file_name,
        "technical": technical_msg,
        "errors": [friendly_msg if user_friendly else technical_msg],
    }


def resize_image(
    image_bytes: bytes,
    max_part_size_bytes: int = 2 * 1024 * 1024,
    strict_target_kb: int = 50,
    max_dimension: int = 480,
    max_attempts: int = 7,
    _logger=None,
) -> Optional[bytes]:
    """
    Aggressively resizes an image to be under a target size.

    It first tries to stay below the `max_part_size_bytes` (4MB) with good quality.
    If the original image is very large, it will then try to aggressively shrink
    it into a "thumbnail" under `strict_target_kb` to help with total payload size.
    Returns None if it cannot shrink the image sufficiently.
    """
    logger = _logger or get_logger()
    strict_target_bytes = strict_target_kb * 1024
    initial_size = len(image_bytes)

    # 1. Quick check: If the image is already thumbnail-sized, we're done.
    if initial_size < strict_target_bytes:
        return image_bytes

    try:
        with PILImage.open(io.BytesIO(image_bytes)) as img:
            img = _normalize_for_jpeg(img)
            output_format = "JPEG"

            # 2. Iterative Resizing Loop
            # start with decent quality and dimensions
            current_quality = 85
            current_max_dim = max_dimension

            for attempt in range(max_attempts):
                # create a copy to avoid re-resizing the same object
                img_copy = img.copy()
                img_copy.thumbnail((current_max_dim, current_max_dim))

                buffer = io.BytesIO()
                img_copy.save(
                    buffer, format=output_format, quality=current_quality, optimize=True
                )
                resized_bytes = buffer.getvalue()

                # a. ideal case: under the strict target. Return immediately.
                if len(resized_bytes) <= strict_target_bytes:

                    return resized_bytes

                # b. good enough case: under the 2MB part limit (only for first attempt)
                if attempt == 0 and len(resized_bytes) <= max_part_size_bytes:
                    logger.debug(
                        f"Image resized to standard target: {len(resized_bytes) / 1024:.1f} KB. "
                        f"Continuing to shrink for strict target."
                    )
                    # continue loop to try and get it even smaller.

                # C. still too big: reduce quality and dimensions for the next loop
                current_quality = max(10, current_quality - 12)
                # drastically reduce dimensions after a few tries
                if attempt > 1:
                    current_max_dim = max(128, int(current_max_dim * 0.75))

            # 3. final check: if loop finishes, check if the last attempt was at least under the 2MB limit
            final_bytes = resized_bytes
            if len(final_bytes) <= max_part_size_bytes:
                logger.debug(
                    f"Could not meet strict target, but image is under part limit. "
                    f"Final size: {len(final_bytes) / 1024:.1f} KB."
                )
                return final_bytes
            else:
                logger.debug(
                    f"Dropping image. Could not shrink below part limit of {max_part_size_bytes / 1024} KB "
                    f"after {max_attempts} attempts. Final size: {len(final_bytes) / 1024 / 1024:.2f} MB."
                )
                return None

    except Exception as e:
        logger.debug(f"Failed to process image for resizing, dropping it. Error: {e}")
        return None


def extract_image_safe(pdf_doc: fitz.Document, xref: int) -> Tuple[bytes, str]:
    """
    Returns (image_bytes, ext) where ext is png | jpeg | ...
    Handles jp2, JPX and CMYK colorspace conversion.
    """
    logger = get_logger()
    base = pdf_doc.extract_image(xref)
    img_bytes, img_ext = base["image"], base["ext"].lower()

    if img_ext == "jpx":
        try:
            pix = fitz.Pixmap(pdf_doc, xref)

            if pix.colorspace.name not in ("GRAY", "RGB"):
                pix = fitz.Pixmap(fitz.csRGB, pix)

            img_bytes = pix.tobytes("png")
            img_ext = "png"

        except Exception as e:
            log_msg = f"Failed to convert JPX to PNG (xref={xref}): {e}"
            logger.debug(log_msg)
            raise RuntimeError(log_msg) from e

    elif img_ext == "jb2":
        try:
            pix = fitz.Pixmap(pdf_doc, xref)
            img_bytes = pix.tobytes("png")
            img_ext = "png"

        except Exception as e:
            log_msg = f"Failed to convert jb2 to PNG (xref={xref}): {e}"
            logger.debug(log_msg)
            raise RuntimeError(log_msg) from e

    return img_bytes, img_ext

def extract_pdf_pass_one(
    pdf_document: fitz.Document, filename: str
) -> Tuple[List[Dict[str, Any]], List[Dict[str, Any]]]:
    """
    High-fidelity extraction. Extracts text and resizes individual images.

    For pages that look like *scanned* pages (little/no extractable text but
    one or more images), we collapse everything into a single full-page
    screenshot instead of keeping multiple layered images.
    """
    logger = get_logger()

    content: List[Dict[str, Any]] = []
    errors: List[Dict[str, Any]] = []
    processed_hashes: set[str] = set()

    # tune this if needed
    SCANNED_TEXT_CHAR_THRESHOLD = 200

    # Early-abort tracking: if pass-one clearly exceeded payload budget,
    # skip expensive extraction for remaining pages and let pass-two handle all pages.
    _cum_content_bytes = 0
    _budget_exceeded = False

    for page_idx in range(len(pdf_document)):
        if _budget_exceeded:
            continue

        page = pdf_document.load_page(page_idx)
        page_no = page_idx + 1

        had_content = False
        had_errors = False

        # 1) TEXT
        text = page.get_text().strip()
        if text:
            content.append(
                {
                    "type": "text",
                    "page": page_no,
                    "content": text,
                    "filename": filename,
                }
            )
            _cum_content_bytes += len(text.encode("utf-8"))
            had_content = True

        # 2) IMAGES
        image_infos = page.get_image_info(xrefs=True) or []

        valid_images: List[Tuple[bytes, str, dict]] = []

        for info in image_infos:
            xref = info.get("xref")
            if not xref:
                continue

            # dedupe WITHOUT hashing bytes
            digest = info.get("digest")
            if isinstance(digest, (bytes, bytearray)):
                img_id = digest.hex()
            elif digest:
                img_id = str(digest)
            else:
                img_id = f"xref:{xref}"

            if img_id in processed_hashes:
                continue
            processed_hashes.add(img_id)

            try:
                img_bytes, img_ext = extract_image_safe(pdf_document, xref)

                if img_ext not in {"png", "jpeg", "jpg", "bmp", "gif", "jb2"}:
                    logger.debug("Page %d: unexpected image format '%s' - skipping", page_no, img_ext)
                    had_errors = True
                    continue

                if not is_valid_image(img_bytes):
                    logger.debug("Page %d: invalid image stream - skipping", page_no)
                    had_errors = True
                    continue

                valid_images.append((img_bytes, img_ext, info))

            except Exception as e:
                tech = f"Page {page_no}: image extract error (xref {xref}): {e}"
                logger.debug(tech)
                errors.append(_wrap_error(filename, tech, user_friendly=True))
                had_errors = True
                continue

        # 2a) SCANNED-PAGE HANDLING
        # If there is very little extractable text but at least one image,
        # treat the page as a scanned image and just keep ONE full-page screenshot.
        countable = sum(1 for info in image_infos if _px_from_info(info) >= _COUNT_MIN_PIXELS)
        is_scanned_like = (len(text) < SCANNED_TEXT_CHAR_THRESHOLD) and (countable > 0)

        if is_scanned_like:
            try:
                page_png = render_pdf_page_as_image(page)  # full quality
                content.append(
                    {
                        "type": "full_page_image",
                        "page": page_no,
                        "content": base64.b64encode(page_png).decode("utf-8"),
                        "format": "png",
                        "filename": filename,
                    }
                )
                _cum_content_bytes += len(page_png)
                had_content = True
                logger.debug(
                    "Page %d flagged as scanned-like (len(text)=%d, images=%d); "
                    "stored single full_page_image.",
                    page_no,
                    len(text),
                    len(valid_images),
                )
                # skip the rest of per-image logic for this page
                if _cum_content_bytes > MAX_INLINE_PAYLOAD_BYTES:
                    _budget_exceeded = True
                    logger.debug(
                        "Page %d: cumulative content ~%dMB exceeds budget; "
                        "skipping extraction for remaining pages.",
                        page_no,
                        _cum_content_bytes // (1024 * 1024),
                    )
                continue
            except Exception as e:
                tech = f"Page {page_no}: full-page render failed (scanned heuristic): {e}"
                logger.debug(tech)
                errors.append(_wrap_error(filename, tech, user_friendly=True))
                # fall through to normal logic if screenshot failed

        # 2b) too-many-images logic
        if image_infos:
            _mark = len(content)
            handled = _collapse_layered_images_to_screenshot_meta(
                page=page,
                filename=filename,
                page_no=page_no,
                image_infos=image_infos,
                content=content,
                errors=errors,
            )
            if handled:
                for ci in range(_mark, len(content)):
                    if content[ci].get("type") != "text":
                        _cum_content_bytes += len(content[ci]["content"]) * 3 // 4
                had_content = True
                if _cum_content_bytes > MAX_INLINE_PAYLOAD_BYTES:
                    _budget_exceeded = True
                    logger.debug(
                        "Page %d: cumulative content ~%dMB exceeds budget; "
                        "skipping extraction for remaining pages.",
                        page_no,
                        _cum_content_bytes // (1024 * 1024),
                    )
                continue


        if countable > MAX_IMAGES_PER_PAGE:
            # Too many countable images -> send ONE full-page screenshot instead
            try:
                page_png = render_pdf_page_as_image(page)  # full quality
                content.append(
                    {
                        "type": "full_page_image",
                        "page": page_no,
                        "content": base64.b64encode(page_png).decode("utf-8"),
                        "format": "png",
                        "filename": filename,
                    }
                )
                _cum_content_bytes += len(page_png)
                had_content = True
            except Exception as e:
                tech = f"Page {page_no}: full-page render failed: {e}"
                logger.debug(tech)
                errors.append(_wrap_error(filename, tech, user_friendly=True))
        else:
            # Under the threshold: keep ALL extracted images at full quality
            for img_bytes, img_ext, info in valid_images:
                fmt = "jpeg" if img_ext.lower() == "jpg" else img_ext.lower()
                content.append(
                    {
                        "type": "image",
                        "page": page_no,
                        "content": base64.b64encode(img_bytes).decode("utf-8"),
                        "format": fmt,
                        "filename": filename,
                    }
                )
                _cum_content_bytes += len(img_bytes)
                had_content = True

        # 3) If still nothing on page, fallback to full-page screenshot
        if not had_content:
            try:
                full_png = render_pdf_page_as_image(page)
                content.append(
                    {
                        "type": "full_page_image",
                        "page": page_no,
                        "content": base64.b64encode(full_png).decode("utf-8"),
                        "format": "png",
                        "filename": filename,
                    }
                )
                _cum_content_bytes += len(full_png)
                had_content = True
            except Exception as e:
                tech = f"Page {page_no}: full-page render failed: {e}"
                logger.debug(tech)
                errors.append(_wrap_error(filename, tech, user_friendly=True))

        if _cum_content_bytes > MAX_INLINE_PAYLOAD_BYTES and not _budget_exceeded:
            _budget_exceeded = True
            logger.debug(
                "Page %d: cumulative content ~%dMB exceeds budget; "
                "skipping extraction for remaining pages.",
                page_no,
                _cum_content_bytes // (1024 * 1024),
            )

    return content, errors


def extract_pdf_pass_two(
    pdf_document: fitz.Document, filename: str
) -> Tuple[List[Dict[str, Any]], List[Dict[str, Any]]]:
    """
    Aggressive "screenshot" extraction for oversized documents.
    Extracts text and a low-res image of every single page.
    """
    logger = get_logger()

    content: List[Dict[str, Any]] = []
    errors: List[Dict[str, Any]] = []

    logger.debug(f"'{filename}' is oversized. Running aggressive screenshot pass.")

    for page_idx in range(len(pdf_document)):
        page = pdf_document.load_page(page_idx)
        page_no = page_idx + 1

        # 1. Always extract text
        text = page.get_text().strip()
        if text:
            content.append(
                {"type": "text", "page": page_no, "content": text, "filename": filename}
            )

        # 2. Always ATTEMPT to take a screenshot of the page
        try:
            page_screenshot_bytes = render_pdf_page_as_image(page)

            capped, _mime = prepare_image_part_bytes(
                page_screenshot_bytes,
                max_part_size_bytes=2 * 1024 * 1024,
                strict_target_kb=200,
                max_dimension=1600,
            )

            if capped:
                content.append(
                        {
                            "type": "full_page_image",
                            "page": page_no,
                            "content": base64.b64encode(capped).decode("utf-8"),
                            "format": "jpeg",
                            "filename": filename,
                        }
                    )
        except Exception as e:
            tech = f"Page {page_no}: Screenshot failed during pass two: {e}"
            logger.debug(tech)
            errors.append(_wrap_error(filename, tech))

    return content, errors


def extract_pdf_text_and_images(pdf_path: str) -> Tuple[List[Dict[str, Any]], List[Dict[str, Any]]]:
    logger = get_logger()
    filename = os.path.basename(pdf_path)

    with _pymupdf_lock:
        with muted_c_stderr():
            try:
                pdf_document = fitz.open(pdf_path)
            except Exception as e:
                tech = f"Error opening PDF {filename}: {e}"
                logger.debug(tech)
                return [], [_wrap_error(filename, tech)]

        try:
            content_pass_one, errors_pass_one = extract_pdf_pass_one(pdf_document, filename)

            total_size = 0
            for item in content_pass_one:
                if item["type"] == "text":
                    total_size += len(item["content"].encode("utf-8"))
                else:
                    total_size += len(base64.b64decode(item["content"]))

            if total_size <= MAX_INLINE_PAYLOAD_BYTES:
                logger.debug(
                    f"'{filename}' is compliant after pass one ({total_size / 1024 / 1024:.2f} MB). Keeping high quality."
                )
                return content_pass_one, errors_pass_one

            logger.debug(
                f"'{filename}' is oversized after pass one ({total_size / 1024 / 1024:.2f} MB). "
                "Switching to aggressive screenshot strategy."
            )
            content_pass_two, errors_pass_two = extract_pdf_pass_two(pdf_document, filename)

            return content_pass_two, errors_pass_one + errors_pass_two

        finally:
            try:
                pdf_document.close()
            except Exception:
                pass


def save_image_to_temp_file(item: dict) -> str:
    """
    Writes the decoded image (from the base64 'content') to a temporary file.
    Returns the file path.
    """
    image_bytes = base64.b64decode(item["content"])
    ext = item["format"]
    temp_file = tempfile.NamedTemporaryFile(delete=False, suffix=f".{ext}")
    temp_file.write(image_bytes)
    temp_file.close()
    return temp_file.name


def render_pdf_page_as_image(page, dpi: int = 150) -> bytes:
    """
    Renders the full PDF page as an image (PNG) using PyMuPDF.
    """
    zoom = dpi / 72.0  # PyMuPDF default is 72 dpi
    mat = fitz.Matrix(zoom, zoom)

    with muted_c_stderr():
        pix = page.get_pixmap(matrix=mat, alpha=False)

    return pix.tobytes("png")


def extract_text_from_docx(filepath: str) -> str:
    doc = docx.Document(filepath)
    text = []

    # Paragraphs
    for p in doc.paragraphs:
        if p.text.strip():
            text.append(p.text)

    # Tables
    for table in doc.tables:
        for row in table.rows:
            row_text = []
            for cell in row.cells:
                cell_text = cell.text.strip()
                if cell_text:
                    row_text.append(cell_text)
            if row_text:
                text.append(" | ".join(row_text))

    return "\n".join(text)


def extract_text_from_txt(filepath: str) -> str:
    """Extracts text from a TXT file."""
    logger = get_logger()

    try:
        with open(filepath, "r", encoding="utf-8", errors="replace") as file:
            return file.read()
    except Exception as e:

        logger.debug(f"Error extracting text from TXT {filepath}: {e}")
        return ""



def extract_text_from_xlsx(filepath: str,
                           *,
                           max_rows_per_sheet: int = 1000,
                           max_cols_per_sheet: int = 100,
                           sheet_cap: int = 30) -> str:
    """Extract readable CSV-like text from ALL sheets in an .xlsx."""

    logger = get_logger()
    out_chunks = []

    def _non_empty(df: pd.DataFrame) -> bool:
        if df is None or df.empty:
            return False
        # treat '','NaN','None' as empty
        mask = df.map(lambda x: (str(x).strip() if x is not None else ""))
        return (mask != "").values.any()

    try:
        xls = pd.ExcelFile(filepath, engine="openpyxl")
    except Exception as e:
        logger.debug(f"ExcelFile(openpyxl) failed, retrying with default engine: {e}")
        xls = pd.ExcelFile(filepath)

    sheets = (xls.sheet_names or [])[:sheet_cap]

    for sheet in sheets:
        df_best = None
        for hdr in (0, None):
            try:
                df = xls.parse(sheet_name=sheet, header=hdr, dtype=str, keep_default_na=False)
                df = df.astype(str)
                df = df.replace({"None": "", "nan": "", "NaN": ""})
                nonempty_cols = df.map(lambda s: s.strip() != "").any(axis=0)
                df = df.loc[:, nonempty_cols]
                nonempty_rows = df.map(lambda s: s.strip() != "").any(axis=1)
                df = df.loc[nonempty_rows, :]
                if _non_empty(df):
                    df_best = df
                    break
            except Exception as e:
                logger.debug(f"parse(sheet={sheet}, header={hdr}) failed: {e}")

        if df_best is None:
            continue

        # clip size
        df_best = df_best.iloc[:max_rows_per_sheet, :max_cols_per_sheet]

        # emit as CSV under a sheet tag
        buf = io.StringIO()
        df_best.to_csv(buf, index=False)
        text_block = buf.getvalue().strip()
        if text_block:
            out_chunks.append(f"<sheet name='{sheet}'>\n{text_block}\n</sheet>")

    # Fallback: try a very defensive openpyxl scan if nothing came out
    if not out_chunks:
        try:
            from openpyxl import load_workbook
            wb = load_workbook(filepath, data_only=True, read_only=True)
            for ws in wb.worksheets[:sheet_cap]:
                rows = []
                for row in ws.iter_rows(values_only=True):
                    row = ["" if v is None else str(v) for v in row]
                    if any(c.strip() for c in row):
                        rows.append(",".join(row))
                        if len(rows) >= max_rows_per_sheet:
                            break
                if rows:
                    out_chunks.append(f"<sheet name='{ws.title}'>\n" + "\n".join(rows) + "\n</sheet>")
        except Exception as e:
            logger.debug(f"openpyxl fallback failed: {e}")

    return "\n\n".join(out_chunks)


def extract_images_from_xlsx(
    filepath: str,
    *,
    sheet_cap: int = 30,
    max_images: int = 80,
    max_images_per_sheet: int = 12,
) -> List[Dict[str, Any]]:
    """
    Extract embedded worksheet images from an .xlsx.
    Returns metadata + raw bytes so caller can normalize via prepare_image_part_bytes().
    """
    logger = get_logger()
    images: List[Dict[str, Any]] = []

    try:
        from openpyxl import load_workbook
        from openpyxl.utils import get_column_letter
        wb = load_workbook(filepath, data_only=True, read_only=False)
    except Exception as e:
        logger.debug(f"load_workbook failed for image extraction ({filepath}): {e}")
        return images

    try:
        for ws in wb.worksheets[:sheet_cap]:
            ws_images = list(getattr(ws, "_images", []) or [])
            if not ws_images:
                continue

            per_sheet = 0
            for img in ws_images:
                if len(images) >= max_images or per_sheet >= max_images_per_sheet:
                    break

                try:
                    raw = img._data()
                    if not raw:
                        continue
                except Exception as e:
                    logger.debug(
                        "Could not read embedded Excel image in sheet %s: %s",
                        ws.title,
                        e,
                    )
                    continue

                row = None
                col = None
                try:
                    anchor = getattr(img, "anchor", None)
                    if anchor is not None and hasattr(anchor, "_from"):
                        row = int(anchor._from.row) + 1
                        col = int(anchor._from.col) + 1
                except Exception:
                    pass

                cell = None
                if row and col:
                    try:
                        cell = f"{get_column_letter(col)}{row}"
                    except Exception:
                        cell = None

                images.append(
                    {
                        "sheet": ws.title,
                        "row": row,
                        "col": col,
                        "cell": cell,
                        "bytes": raw,
                    }
                )
                per_sheet += 1

            if len(images) >= max_images:
                logger.debug(
                    "Reached Excel image cap (%d) while processing %s",
                    max_images,
                    os.path.basename(filepath),
                )
                break

    finally:
        try:
            wb.close()
        except Exception:
            pass

    return images


def extract_text_from_doc(file_path: str) -> str:
    """
    Extracts text from a legacy .doc (binary Word) file.

    Strategy:
      1) antiword
      2) pypandoc
      3) python-docx (for mislabeled .docx)
    """
    logger = get_logger()

    # 1) antiword for true binary .doc files
    try:
        import subprocess

        result = subprocess.run(
            ["antiword", file_path], capture_output=True, text=True, timeout=60
        )
        if result.returncode == 0 and result.stdout.strip():
            logger.debug(
                "Extracted DOC via antiword: %s (%d chars)",
                os.path.basename(file_path),
                len(result.stdout),
            )
            return result.stdout
    except FileNotFoundError:
        logger.debug("antiword not installed; trying fallbacks for %s", file_path)
    except Exception as e:
        logger.debug("antiword failed for %s: %s", file_path, e)

    # 2) pypandoc fallback
    try:
        text = pypandoc.convert_file(file_path, "plain")
        if text.strip():
            return text
    except Exception as e:
        logger.debug(f"pypandoc failed for DOC {file_path}: {e}")

    # 3) python-docx fallback (some .doc files are actually .docx)
    try:
        from docx import Document

        doc = Document(file_path)
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        if paragraphs:
            text = "\n".join(paragraphs)
            logger.debug(
                "Extracted DOC via python-docx (likely .docx renamed): %s",
                os.path.basename(file_path),
            )
            return text
    except Exception as e:
        logger.debug("python-docx failed for DOC %s: %s", file_path, e)

    logger.warning(
        "Could not extract text from DOC %s (tried antiword, pypandoc, python-docx).",
        os.path.basename(file_path),
    )
    return ""


def is_valid_image(image_bytes: bytes) -> bool:
    """Checks if image bytes represent a valid image."""
    logger = get_logger()

    try:
        # logger.debug(f"      is_valid_image: Attempting to open image ({len(image_bytes)} bytes)...")  # Debugging
        image = PILImage.open(io.BytesIO(image_bytes))
        # logger.debug(f"      is_valid_image: Image opened successfully. Format: {image.format}") # crucial
        return True
    except Exception as e:

        logger.debug(f"      is_valid_image: Failed to open image: {e}")
        return False

def extract_text_from_pptx(filepath: str) -> List[Tuple[int, str]]:
    """
    Returns list of (slide_number, text) with all text found on each slide.
    Includes text frames + tables + grouped shapes. Ignores media.
    """
    logger = get_logger()
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
    except Exception as e:
        logger.debug(f"python-pptx not installed or import failed: {e}")
        raise

    prs = Presentation(filepath)

    def iter_shapes(shapes):
        for shp in shapes:
            yield shp
            # recurse groups
            try:
                if shp.shape_type == MSO_SHAPE_TYPE.GROUP and hasattr(shp, "shapes"):
                    yield from iter_shapes(shp.shapes)
            except Exception:
                pass

    out: List[Tuple[int, str]] = []
    for si, slide in enumerate(prs.slides, start=1):
        chunks: List[str] = []

        for shape in iter_shapes(slide.shapes):
            # Text boxes / titles / placeholders / etc.
            try:
                if getattr(shape, "has_text_frame", False) and shape.text_frame:
                    txt = shape.text_frame.text
                    if txt and txt.strip():
                        chunks.append(txt.strip())
            except Exception:
                pass

            # Tables
            try:
                if getattr(shape, "has_table", False) and shape.table:
                    for row in shape.table.rows:
                        row_cells = []
                        for cell in row.cells:
                            t = (cell.text or "").strip()
                            if t:
                                row_cells.append(t)
                        if row_cells:
                            chunks.append(" | ".join(row_cells))
            except Exception:
                pass

        slide_text = "\n".join(chunks).strip()
        if slide_text:
            out.append((si, slide_text))

    return out

def process_document_parts(
    doc_path: str, doc_name: str, index: int
) -> Tuple[List[Part], List[Dict[str, Any]], List[dict]]:
    """Processes a document (PDF, DOCX, DOC, TXT, XLSX, or image) into parts."""
    logger = get_logger()

    parts: List[Part] = []
    errors: List[Dict[str, Any]] = []
    tokenized_pages: List[dict] = []

    if not os.path.exists(doc_path):
        error_msg = f"Document {index}: {doc_name} - File not found."
        logger.debug(error_msg)
        errors.append(_wrap_error(doc_name, "File not found.", user_friendly=True))
        parts.append(Part.from_text(text=error_msg))
        return parts, errors, []

    ext = os.path.splitext(doc_name)[1].lower()
    open_tag = f"<Start of File: {doc_name}>\n"
    parts.append(Part.from_text(text=open_tag))

    try:
        # Handle archives
        if ext in ARCHIVES:
            # Stage into a private working dir to keep contents isolated
            work_root = tempfile.mkdtemp(prefix=f"bundle_{Path(doc_name).stem}_")
            try:
                staged_archive = os.path.join(work_root, doc_name)
                shutil.copy2(doc_path, staged_archive)

                zip_results = process_zip_files_in_directory(
                    work_root, root_directory=work_root
                )
                if zip_results:
                    # normalize error format a bit
                    for r in zip_results:
                        errors.append(
                            _wrap_error(
                                r.get("file", doc_name),
                                "; ".join(
                                    {e.get("error", "") for e in r.get("errors", [])}
                                )
                                or "Archive processing issue",
                                user_friendly=True,
                            )
                        )

                # collect all supported files that resulted from extraction
                inner_files: List[str] = []
                for root, _, files in os.walk(work_root):
                    for fn in files:
                        if fn.startswith(".") or fn in {"__MACOSX", ".DS_Store"}:
                            continue
                        if (
                            Path(fn).suffix.lower() in SUPPORTED_DOCS
                            or Path(fn).suffix.lower() in ARCHIVES
                        ):
                            inner_files.append(os.path.join(root, fn))

                # deterministic order for reproducibility
                inner_files.sort(key=lambda p: p.lower())

                # recursively process every inner item (archives included)
                for i, ipath in enumerate(inner_files):
                    iname = os.path.basename(ipath)
                    rel = os.path.relpath(ipath, work_root)
                    sub_parts, sub_errors, sub_tok = process_document_parts(
                        ipath, iname, i
                    )
                    parts.extend(sub_parts)
                    errors.extend(sub_errors or [])
                    # tag tokenized pages with the inner file for disambiguation
                    for t in sub_tok:
                        if "file" not in t:
                            t["file"] = iname
                        # Compose a stable relative path (handles nested archives too)
                        if "file_path" in t:
                            t["file_path"] = os.path.join(rel, t["file_path"])
                        else:
                            t["file_path"] = rel
                        tokenized_pages.append(t)

            finally:
                # clean up the staging area
                try:
                    shutil.rmtree(work_root)
                except Exception as ex:
                    logger.debug(f"Could not clean temp bundle dir {work_root}: {ex}")

        elif ext in [".png", ".jpg", ".jpeg", ".gif", ".bmp", ".webp"]:
            try:
                with open(doc_path, "rb") as imgf:
                    raw = imgf.read()

                capped, mime_type = prepare_image_part_bytes(
                    raw,
                    max_part_size_bytes=2 * 1024 * 1024,  # SAME cap
                    strict_target_kb=200,
                    max_dimension=1600,
                )
                if not capped:
                    raise ValueError("Image could not be shrunk under the per-part size cap")

                parts.append(types.Part.from_bytes(data=capped, mime_type=mime_type))

            except Exception as ex:
                err = f"[Error reading/resizing image {doc_name}: {ex}]"
                logger.debug(err)
                errors.append(_wrap_error(doc_name, str(ex), user_friendly=True))
                parts.append(Part.from_text(text=f"[Image {doc_name} omitted: {ex}]"))

        elif ext == ".pptx":
            try:
                slides = extract_text_from_pptx(doc_path)
                if not slides:
                    raise ValueError("No text extracted; PPTX may be empty or text is not accessible")

                # Add every slide text as its own text Part
                for slide_no, slide_text in slides:
                    parts.append(
                        Part.from_text(
                            text=f"Slide {slide_no} of {doc_name}:\n{slide_text}"
                        )
                    )
                    tokenized_pages.append(
                        {"type": "text", "page": slide_no, "content": slide_text, "file": doc_name}
                    )

            except Exception as ex:
                error_msg = f"Error extracting PPTX from {doc_name}: {ex}"
                logger.debug(error_msg)
                errors.append(_wrap_error(doc_name, error_msg, user_friendly=True))
                parts.append(Part.from_text(text=error_msg))

        elif ext == ".pdf":
            pdf_items, pdf_errors = extract_pdf_text_and_images(doc_path)
            pdf_items.sort(key=lambda item: item["page"])
            if pdf_errors:
                for err in pdf_errors:
                    logger.debug("Error in PDF (%s): %s", err["file"], err["technical"])
                errors.extend(pdf_errors)

            # Pages that already have text extracted (i.e., not pure scanned images)
            pages_with_text = {
                item["page"]
                for item in pdf_items
                if item.get("type") == "text" and str(item.get("content", "")).strip()
            }

            # Process image items in parallel while preserving output order.
            image_work: List[Tuple[int, dict]] = []
            for i, item in enumerate(pdf_items):
                if item["type"] in ("image", "full_page_image"):
                    image_work.append((i, item))

            image_results: Dict[int, Tuple[List[Part], List[dict], List[dict]]] = {}
            _parent_logger = logger

            def _process_one_pdf_image(
                idx: int, item: dict, _pages_with_text: set, _doc_name: str
            ) -> Tuple[int, List[Part], List[dict], List[dict]]:
                img_parts: List[Part] = []
                img_tokenized: List[dict] = []
                img_errors: List[dict] = []

                image_bytes = base64.b64decode(item["content"])
                try:
                    image_for_model, mime_type = prepare_image_part_bytes(
                        image_bytes, _logger=_parent_logger
                    )
                    if not image_for_model or not mime_type:
                        raise RuntimeError("image could not be normalized for model input")

                    from utils.document.ocr_deepseek import ENABLE_IMAGE_OCR, ocr_image_to_text

                    page_no = item.get("page")
                    page_has_text = page_no in _pages_with_text
                    try:
                        pixels = _image_pixel_count(image_for_model)
                    except Exception:
                        pixels = 0
                    is_big = pixels >= 400_000

                    ocr_text = ""
                    if ENABLE_IMAGE_OCR and (not page_has_text) and is_big:
                        try:
                            ocr_text = ocr_image_to_text(
                                image_for_model,
                                mime_type=mime_type,
                                filename_hint=f"{_doc_name}_p{page_no}",
                            )
                        except Exception as ocr_ex:
                            _parent_logger.debug(
                                "OCR skipped/failed for %s page %s: %s",
                                _doc_name,
                                page_no,
                                ocr_ex,
                            )
                            ocr_text = ""

                    if ocr_text:
                        img_parts.append(
                            Part.from_text(
                                text=f"[Extracted text from image on page {page_no} of {_doc_name}]\n{ocr_text}"
                            )
                        )
                        img_tokenized.append(
                            {
                                "type": "ocr_text",
                                "page": page_no,
                                "content": ocr_text,
                                "file": _doc_name,
                            }
                        )

                    img_parts.append(
                        types.Part.from_bytes(data=image_for_model, mime_type=mime_type)
                    )

                except Exception as e:
                    error_msg = (
                        f"Error converting image {item.get('filename', 'unknown')} "
                        f"on page {item.get('page', '?')}: {e}"
                    )
                    _parent_logger.debug(error_msg)
                    img_errors.append(
                        _wrap_error(
                            item.get("filename", "unknown"),
                            error_msg,
                            user_friendly=True,
                        )
                    )

                return idx, img_parts, img_tokenized, img_errors

            if image_work:
                with ThreadPoolExecutor(max_workers=_IMAGE_PROCESS_WORKERS) as img_pool:
                    futures = {
                        img_pool.submit(
                            _process_one_pdf_image, idx, item, pages_with_text, doc_name
                        ): idx
                        for idx, item in image_work
                    }
                    for fut in as_completed(futures):
                        try:
                            idx, img_p, img_t, img_e = fut.result()
                            image_results[idx] = (img_p, img_t, img_e)
                        except Exception as e:
                            logger.debug("Image processing thread failed: %s", e)

            # Reassemble in original page order.
            for i, item in enumerate(pdf_items):
                if item["type"] == "text":
                    parts.append(
                        Part.from_text(
                            text=f"Page {item['page']} of {item.get('filename', 'unknown')}:\n{item['content']}"
                        )
                    )
                    tokenized_pages.append(
                        {
                            "type": "text",
                            "page": item["page"],
                            "content": item["content"],
                            "file": doc_name,
                        }
                    )
                elif i in image_results:
                    img_p, img_t, img_e = image_results[i]
                    parts.extend(img_p)
                    tokenized_pages.extend(img_t)
                    errors.extend(img_e)

        elif ext == ".docx":
            try:
                text = extract_text_from_docx(doc_path)
                if not text.strip():
                    raise ValueError("No text extracted; possibly corrupt DOCX")
                parts.append(Part.from_text(text=text))
                tokenized_pages.append(
                    {"type": "text", "page": 1, "content": text, "file": doc_name}
                )

            except Exception as ex:
                error_msg = f"Error extracting DOCX from {doc_name}: {ex}"
                logger.debug(error_msg)
                errors.append(_wrap_error(doc_name, error_msg, user_friendly=True))
                parts.append(Part.from_text(text=error_msg))

        elif ext == ".doc":
            try:
                text = extract_text_from_doc(doc_path)
                if not text.strip():
                    raise ValueError("No text extracted; possibly corrupt DOC")
                parts.append(Part.from_text(text=text))
                tokenized_pages.append(
                    {"type": "text", "page": 1, "content": text, "file": doc_name}
                )

            except Exception as ex:
                error_msg = f"Error extracting DOC from {doc_name}: {ex}"
                logger.debug(error_msg)
                errors.append(_wrap_error(doc_name, error_msg, user_friendly=True))
                parts.append(Part.from_text(text=error_msg))

        elif ext == ".txt":
            try:
                text = extract_text_from_txt(doc_path)
                if not text.strip():
                    raise ValueError("No text extracted; possibly corrupt TXT")
                parts.append(Part.from_text(text=text))
                tokenized_pages.append(
                    {"type": "text", "page": 1, "content": text, "file": doc_name}
                )

            except Exception as ex:
                error_msg = f"Error extracting TXT from {doc_name}: {ex}"
                logger.debug(error_msg)
                errors.append(_wrap_error(doc_name, error_msg, user_friendly=True))
                parts.append(Part.from_text(text=error_msg))

        elif ext == ".xlsx":
            try:
                text = extract_text_from_xlsx(doc_path)
                if not text.strip():
                    raise ValueError("No text extracted; possibly corrupt xlsx")
                parts.append(Part.from_text(text=text))
                tokenized_pages.append(
                    {"type": "text", "page": 1, "content": text, "file": doc_name}
                )

                embedded_images = extract_images_from_xlsx(doc_path)
                for img in embedded_images:
                    image_bytes = img.get("bytes")
                    if not image_bytes:
                        continue

                    sheet_name = str(img.get("sheet") or "unknown")
                    cell = img.get("cell")
                    where = f"sheet '{sheet_name}'"
                    if cell:
                        where = f"{where}, cell {cell}"

                    image_for_model, mime_type = prepare_image_part_bytes(image_bytes)
                    if not image_for_model:
                        msg = f"Embedded Excel image could not be normalized ({where})"
                        logger.debug("%s in %s", msg, doc_name)
                        errors.append(_wrap_error(doc_name, msg, user_friendly=True))
                        continue

                    parts.append(
                        Part.from_text(text=f"[Embedded image from {where} of {doc_name}]")
                    )
                    parts.append(
                        types.Part.from_bytes(data=image_for_model, mime_type=mime_type)
                    )
                    tokenized_pages.append(
                        {
                            "type": "image",
                            "page": 1,
                            "content": f"Embedded image from {where}",
                            "file": doc_name,
                            "sheet": sheet_name,
                            "cell": cell,
                        }
                    )

            except Exception as ex:
                error_msg = f"Error extracting XLSX from {doc_name}: {ex}"
                logger.debug(error_msg)
                errors.append(_wrap_error(doc_name, error_msg, user_friendly=True))
                parts.append(Part.from_text(text=error_msg))
        else:
            error_text = f"Unsupported document type: {doc_name}"
            parts.append(Part.from_text(text=error_text))
            errors.append(
                _wrap_error(doc_name, "Unsupported document type.", user_friendly=True)
            )

    except Exception as e:
        error_msg = f"Error processing document {doc_name}: {e}"
        logger.debug(error_msg)
        errors.append(_wrap_error(doc_name, error_msg, user_friendly=True))
        parts.append(Part.from_text(text=error_msg))

    close_tag = f"</End of File: {doc_name}>\n"
    parts.append(Part.from_text(text=close_tag))
    return parts, errors, tokenized_pages


def process_zip_files_in_directory(directory_path, root_directory=None):
    """
    Recursively process a directory to:
    1. Find and extract zip + 7z files
    2. Remove unsupported files
    3. Process subdirectories
    4. Move supported documents to root directory

    Args:
        directory_path: Path to the directory to process
        root_directory: The top-level directory to move files to (set on first call)

    Returns:
        list: List of errors with unsupported or corrupted files in JSON format
    """
    logger = get_logger()

    logger.debug(f"STARTING to process directory: {directory_path}")
    logger.debug(f"Current root_directory: {root_directory}")
    supported_documents = [
        ".pdf",
        ".docx",
        ".doc",
        ".txt",
        ".xlsx",
        ".png",
        ".jpg",
        ".jpeg",
        ".gif",
        ".bmp",
        ".pptx",
    ]
    results = []

    # on first call, set root_directory to directory_path
    if root_directory is None:
        root_directory = directory_path
        logger.debug(f"Setting initial root_directory to: {root_directory}")

    # ensure directory exists
    if not os.path.exists(directory_path):
        logger.debug(f"Directory does not exist: {directory_path}")
        return results

    logger.debug(f"Processing directory: {directory_path}")

    # First pass: extract all zip files
    for item in os.listdir(directory_path):
        item_path = os.path.join(directory_path, item)

        if os.path.isfile(item_path) and item.lower().endswith(".7z"):
            extract_dir = os.path.splitext(item_path)[0]
            logger.debug(f"Extracting 7z file: {item_path} to {extract_dir}")

            try:
                with py7zr.SevenZipFile(item_path, mode="r") as z:
                    # Extract all contents to the current directory
                    z.extractall(extract_dir)
                    logger.debug(f"Extraction successful to: {extract_dir}")
            except:
                logger.debug(f"Bad 7z file: {item_path}")
                results.append(
                    {
                        "file": os.path.basename(item_path),
                        "errors": [
                            {
                                "doc": os.path.basename(item_path),
                                "error": "Bad 7z file",
                            }
                        ],
                    }
                )
        if os.path.isfile(item_path) and item.lower().endswith(".zip"):
            # logger.debug(f"Found zip file: {item_path}")
            extract_dir = os.path.splitext(item_path)[0]
            logger.debug(f"Extracting zip file: {item_path} to {extract_dir}")

            try:
                with zipfile.ZipFile(item_path, "r") as zip_ref:
                    # logger.debug(f"Zip contains: {zip_ref.namelist()}")
                    zip_ref.extractall(extract_dir)
                    logger.debug(f"Extraction successful to: {extract_dir}")

            except zipfile.BadZipFile:
                # logger.debug(f"ERROR: Bad zip file: {item_path}")
                logger.debug(f"Bad zip file: {item_path}")
                results.append(
                    {
                        "file": os.path.basename(item_path),
                        "errors": [
                            {
                                "doc": os.path.basename(item_path),
                                "error": "Bad zip file",
                            }
                        ],
                    }
                )
            except Exception as e:
                logger.debug(f"Error extracting {item_path}: {str(e)}")
                results.append(
                    {
                        "file": os.path.basename(item_path),
                        "errors": [
                            {
                                "doc": os.path.basename(item_path),
                                "error": f"Extraction error: {str(e)}",
                            }
                        ],
                    }
                )
    # Second pass: process all items (including newly extracted ones)
    # Keep track of zip files to delete them at the end
    zip_files_to_delete = []
    sevenz_files_to_delete = []

    for item in os.listdir(directory_path):
        item_path = os.path.join(directory_path, item)
        # Handle macOS metadata and hidden files
        if (
            (item and item[0] == ".")
            or item == "__MACOSX"
            or item == ".DS_Store"
            or item[0] == "~"
        ):
            logger.debug(f"Removing system file/directory: {item_path}")
            try:
                if os.path.isdir(item_path):
                    shutil.rmtree(item_path)
                else:
                    os.remove(item_path)
                continue
            except Exception as e:
                logger.debug(f"Error removing {item_path}: {str(e)}")
                continue

        if os.path.isdir(item_path):
            # Recursively process subdirectories
            subdirectory_results = process_zip_files_in_directory(
                item_path, root_directory
            )
            results.extend(subdirectory_results)
        elif os.path.isfile(item_path):
            if item.lower().endswith(".zip"):
                zip_files_to_delete.append(item_path)
                continue
            if item.lower().endswith(".7z"):
                sevenz_files_to_delete.append(item_path)
                continue

            file_ext = os.path.splitext(item_path)[1].lower()
            if file_ext not in supported_documents:
                logger.debug(f"Removing unsupported file: {item_path}")
                try:
                    os.remove(item_path)
                    logger.debug(f"Successfully removed unsupported file: {item_path}")
                    results.append(
                        {
                            "file": os.path.basename(item_path),
                            "errors": [
                                {
                                    "doc": os.path.basename(item_path),
                                    "error": "File type not supported",
                                }
                            ],
                        }
                    )
                except Exception as e:
                    logger.debug(
                        f"ERROR: Failed to remove unsupported file {item_path}: {str(e)}"
                    )
                    logger.debug(f"Error removing file {item_path}: {str(e)}")
            else:
                # If this is a supported file AND it's in a subdirectory, move it to root
                if directory_path != root_directory:
                    dest_file = os.path.join(root_directory, item)

                    # Handle filename conflicts by adding a suffix
                    original_name = item
                    if os.path.exists(dest_file):
                        base_name, ext = os.path.splitext(item)
                        counter = 1
                        while os.path.exists(
                            os.path.join(root_directory, f"{base_name}_{counter}{ext}")
                        ):
                            counter += 1
                        item = f"{base_name}_{counter}{ext}"
                        dest_file = os.path.join(root_directory, item)
                        logger.debug(f"Renamed to: {item}")

                    logger.debug(f"Moving file: {item_path} to {dest_file}")
                    try:
                        shutil.move(item_path, dest_file)
                    except Exception as e:
                        logger.debug(f"Error moving file {item_path}: {str(e)}")
                else:
                    logger.debug(
                        f"Keeping supported file in place (already in root): {item_path}"
                    )

    # After all processing is complete, delete the zip files
    for zip_file in zip_files_to_delete:
        try:
            os.remove(zip_file)
            logger.debug(f"Deleted zip file: {zip_file}")
        except Exception as e:
            logger.debug(f"Error deleting zip file {zip_file}: {str(e)}")

    for sevenz_file in sevenz_files_to_delete:
        try:
            os.remove(sevenz_file)
            logger.debug(f"Deleted 7z file: {sevenz_file}")
        except Exception as e:
            logger.debug(f"Error deleting 7z file {sevenz_file}: {str(e)}")

    logger.debug(f"Found {len(results)} issues in this directory")

    return results


def _sanitize_parts_keep_media(parts: list[Part]) -> list[Part]:
    """
    Keep non-empty text Parts and any media/file Parts (inline_data or file_data),
    preserving order. If there is no text at all, append a small fallback text Part.
    """
    cleaned: list[Part] = []
    has_text = False

    for p in parts or []:
        t = getattr(p, "text", None)
        if isinstance(t, str):
            if t.strip():
                cleaned.append(p)
                has_text = True
            # else, drop empty text parts
        else:
            # Keep media/file parts (inline_data or file_data)
            cleaned.append(p)

    if not has_text:
        cleaned.append(
            Part.from_text(
                text="[No extractable text; proceeding with evaluation prompt]"
            )
        )

    return cleaned


def main(*, package_id: str = "system_check") -> bool:
    """
    Lightweight self-test for this module (PDF text/images, ZIP/7z, sanitize,
    resize_image, save_image_to_temp_file, pass-two forcing).
    Prints DOC OK / DOC ERROR:<reason> and returns True/False.
    """
    import base64, tempfile, shutil
    from pathlib import Path

    print("DOC TEST START")
    try:
        sandbox = Path(tempfile.mkdtemp(prefix="doc_selftest_"))
        proj_root = sandbox / package_id
        doc_dir   = proj_root / "documents"
        doc_dir.mkdir(parents=True, exist_ok=True)

        set_logger(pid_tool_logger(package_id, "doc_selftest"))

        # TXT
        (doc_dir / "sample.txt").write_text("This is a sample TXT file.\nSecond line.")

        # DOCX
        d = docx.Document()
        d.add_heading("Sample DOCX", level=1)
        d.add_paragraph("This is some text in the document.")
        d.save(doc_dir / "sample.docx")

        # XLSX
        pd.DataFrame({"Name": ["Alice", "Bob"], "Age": [25, 30]}).to_excel(
            doc_dir / "sample.xlsx", index=False
        )

        # JPG (large-ish)
        big_img_path = doc_dir / "big.jpg"
        PILImage.new("RGB", (2200, 2200), color=(200, 30, 30)).save(
            big_img_path, "JPEG", quality=95
        )

        # PNG
        png_path = doc_dir / "sample.png"
        PILImage.new("RGB", (300, 200), color=(30, 200, 30)).save(png_path, "PNG")

        # PDF with text + embedded image
        pdf_path = doc_dir / "sample.pdf"
        pdf = fitz.open()
        page = pdf.new_page()
        page.insert_text((72, 72), "Sample PDF file content here.")
        page.insert_image(fitz.Rect(72, 120, 300, 300), filename=str(png_path))
        pdf.save(pdf_path)
        pdf.close()

        # BIG-image PDF to force pass-two later
        pdf_big_path = doc_dir / "big_image.pdf"
        pdf2 = fitz.open()
        p2 = pdf2.new_page()
        p2.insert_text((72, 72), "This PDF has a very large embedded image.")
        p2.insert_image(fitz.Rect(72, 100, 1000, 1000), filename=str(big_img_path))
        pdf2.save(pdf_big_path)
        pdf2.close()

        # Unsupported file inside archives to trigger errors list
        unsupported = doc_dir / "note.unsupported"
        unsupported.write_text("not supported type")

        # ZIP archive with mixed contents + nested dir
        zip_root = doc_dir / "bundle"
        zip_root.mkdir(exist_ok=True)
        (zip_root / "inner.txt").write_text("hello from inner txt")
        PILImage.new("RGB", (120, 80), color=(80, 80, 200)).save(zip_root / "inner.jpg", "JPEG")
        (zip_root / "bad.unsupported").write_text("nope")
        nested = zip_root / "nested"
        nested.mkdir(exist_ok=True)
        (nested / "nested.txt").write_text("nested file")
        PILImage.new("RGB", (50, 50), color=(80, 200, 80)).save(nested / "nested.png", "PNG")

        zip_path = doc_dir / "archive.zip"
        with zipfile.ZipFile(zip_path, "w", zipfile.ZIP_DEFLATED) as z:
            for p in [zip_root / "inner.txt", zip_root / "inner.jpg", zip_root / "bad.unsupported",
                      nested / "nested.txt", nested / "nested.png"]:
                z.write(p, arcname=p.relative_to(zip_root))

        # 7z (best-effort; skip silently if not writable)
        sevenz_path = doc_dir / "archive.7z"
        try:
            with py7zr.SevenZipFile(sevenz_path, "w") as z7:
                z7.write(zip_root / "inner.txt", "inner.txt")
                z7.write(zip_root / "inner.jpg", "inner.jpg")
                z7.write(zip_root / "bad.unsupported", "bad.unsupported")
                z7.write(nested / "nested.txt", "nested/nested.txt")
                z7.write(nested / "nested.png", "nested/nested.png")
        except Exception:
            if sevenz_path.exists():
                sevenz_path.unlink()

        # resize_image: force a tiny strict target to make sure we compress
        with open(big_img_path, "rb") as f:
            big_bytes = f.read()
        initial_size = len(big_bytes)
        resized = resize_image(big_bytes, strict_target_kb=5)  # ~5KB target
        if not resized:
            raise RuntimeError("resize_image returned None")
        if len(resized) > initial_size:
            raise RuntimeError("resize_image produced larger output than input")

        # save_image_to_temp_file
        temp_saved = save_image_to_temp_file({
            "content": base64.b64encode(resized).decode("utf-8"),
            "format": "jpeg"
        })
        if not os.path.exists(temp_saved):
            raise RuntimeError("save_image_to_temp_file did not create a file")
        try:
            os.remove(temp_saved)
        except Exception:
            pass

        # extract_pdf_text_and_images (pass one)
        c1, e1 = extract_pdf_text_and_images(str(pdf_path))
        if not isinstance(c1, list) or not all(isinstance(x, dict) for x in c1):
            raise RuntimeError("PDF pass-one extraction returned unexpected format")

        # extract_image_safe sanity on embedded image
        pdf_doc = fitz.open(str(pdf_path))
        try:
            imgs = pdf_doc[0].get_images(full=True)
            if imgs:
                xref = imgs[0][0]
                img_bytes, img_ext = extract_image_safe(pdf_doc, xref)
                if not img_bytes or not isinstance(img_ext, str):
                    raise RuntimeError("extract_image_safe failed on embedded image")
        finally:
            pdf_doc.close()

        # Force pass-two by temporarily shrinking the inline payload limit
        global MAX_INLINE_PAYLOAD_BYTES
        old_limit = MAX_INLINE_PAYLOAD_BYTES
        try:
            MAX_INLINE_PAYLOAD_BYTES = 1024  # 1KB -> forces pass-two path
            c2, e2 = extract_pdf_text_and_images(str(pdf_big_path))
            if not isinstance(c2, list) or not all(isinstance(x, dict) for x in c2):
                raise RuntimeError("PDF pass-two extraction returned unexpected format")
        finally:
            MAX_INLINE_PAYLOAD_BYTES = old_limit

        # process_zip_files_in_directory
        work_dir = doc_dir / "work_archives"
        work_dir.mkdir(exist_ok=True)
        shutil.copy2(zip_path, work_dir / zip_path.name)
        if sevenz_path.exists():
            shutil.copy2(sevenz_path, work_dir / sevenz_path.name)
        z_results = process_zip_files_in_directory(str(work_dir))
        if not isinstance(z_results, list):
            raise RuntimeError("process_zip_files_in_directory did not return a list")

        # process_document_parts on each file
        to_test = [
            ("sample.txt", "TXT"),
            ("sample.docx", "DOCX"),
            ("sample.xlsx", "XLSX"),
            ("sample.png", "PNG"),
            ("sample.pdf", "PDF"),
            ("archive.zip", "ZIP"),
        ]
        if sevenz_path.exists():
            to_test.append(("archive.7z", "7Z"))

        for idx, (name, label) in enumerate(to_test, 1):
            parts, errs, toks = process_document_parts(str(doc_dir / name), name, idx)
            if not isinstance(parts, list) or not all(isinstance(p, Part) for p in parts):
                raise RuntimeError(f"process_document_parts({label}) produced invalid Parts")
            if not isinstance(errs, list) or not isinstance(toks, list):
                raise RuntimeError(f"process_document_parts({label}) returned non-list meta")

            # sanitize & size
            cleaned = _sanitize_parts_keep_media(parts)
            total_sz = estimate_total_parts_size(cleaned)
            if not isinstance(total_sz, int) or total_sz <= 0:
                raise RuntimeError(f"sanitize/size failed for {label}")

        # sanitize fallback when *no* text present (only media/whitespace)
        junk = [
            Part.from_text(text=""),
            Part.from_text(text="  "),
            types.Part.from_bytes(data=b"\x00\x01\x02", mime_type="application/octet-stream"),
        ]
        cleaned = _sanitize_parts_keep_media(junk)
        if not any((getattr(p, "text", None) or "").strip() for p in cleaned):
            raise RuntimeError("sanitize did not insert fallback text when no text present")

        print("DOC OK")
        return True

    except Exception as e:
        print(f"DOC ERROR: {e}")
        return False

    finally:
        try:
            shutil.rmtree(sandbox, ignore_errors=True)  # type: ignore[name-defined]
        except Exception:
            pass


if __name__ == "__main__":
    main()
